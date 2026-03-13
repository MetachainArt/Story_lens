from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent.parent
SCREEN_DIR = BASE_DIR / "assets" / "screens"
OUTPUT_PATH = Path(__file__).resolve().parent / "story-lens-public-proposal.pptx"


class Palette:
    bg = RGBColor(0xF7, 0xF3, 0xEE)
    surface = RGBColor(0xFF, 0xFB, 0xF6)
    line = RGBColor(0xE4, 0xD8, 0xCA)
    ink = RGBColor(0x2E, 0x26, 0x1E)
    muted = RGBColor(0x65, 0x58, 0x4C)
    accent = RGBColor(0xC4, 0x75, 0x50)
    accent_dark = RGBColor(0xA8, 0x5D, 0x3A)
    sage = RGBColor(0x86, 0xA1, 0x89)
    sage_soft = RGBColor(0xE4, 0xEF, 0xE8)
    blue = RGBColor(0x8F, 0xA8, 0xB6)
    blue_soft = RGBColor(0xE7, 0xEF, 0xF5)
    peach_soft = RGBColor(0xF5, 0xE4, 0xD5)
    gold = RGBColor(0xD8, 0xB3, 0x6A)
    white = RGBColor(0xFF, 0xFF, 0xFF)


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = Palette.ink,
    font_name: str = "Malgun Gothic",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    italic: bool = False,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.name = font_name
    font.color.rgb = color


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    *,
    font_size: int = 17,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.word_wrap = True
    for idx, item in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Malgun Gothic"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = Palette.muted
        paragraph.bullet = True


def add_round_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill_color: RGBColor,
    line_color: RGBColor = Palette.line,
    radius_text: str | None = None,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    if radius_text:
        add_textbox(
            slide,
            left + 0.15,
            top + 0.1,
            width - 0.3,
            height - 0.2,
            radius_text,
            font_size=12,
            bold=True,
            color=Palette.accent_dark,
        )


def add_soft_motif(slide) -> None:
    # Top-right friendly bubbles
    for left, top, size, color in [
        (11.5, 0.4, 0.55, Palette.peach_soft),
        (12.15, 0.75, 0.35, Palette.sage_soft),
        (11.85, 1.15, 0.28, Palette.blue_soft),
    ]:
        bubble = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left),
            Inches(top),
            Inches(size),
            Inches(size),
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = color
        bubble.line.color.rgb = color

    # Bottom-left sticker icons
    for i, (label, fill_color) in enumerate(
        [
            ("📷", Palette.peach_soft),
            ("✍", Palette.sage_soft),
            ("♪", Palette.blue_soft),
            ("📘", Palette.peach_soft),
        ]
    ):
        left = 0.45 + i * 0.55
        sticker = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left),
            Inches(6.82),
            Inches(0.42),
            Inches(0.42),
        )
        sticker.fill.solid()
        sticker.fill.fore_color.rgb = fill_color
        sticker.line.color.rgb = Palette.line
        add_textbox(
            slide,
            left,
            6.84,
            0.42,
            0.22,
            label,
            font_size=14,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )


def add_header(
    slide, title: str, subtitle: str | None = None, page_no: int | None = None
) -> None:
    add_textbox(
        slide,
        0.65,
        0.42,
        3.2,
        0.32,
        "Story Lens Proposal",
        font_size=12,
        bold=True,
        color=Palette.accent_dark,
    )
    add_textbox(
        slide,
        0.65,
        0.82,
        7.8,
        0.7,
        title,
        font_size=28,
        bold=True,
        color=Palette.ink,
        font_name="Malgun Gothic",
    )
    if subtitle:
        add_textbox(
            slide, 0.68, 1.42, 8.8, 0.42, subtitle, font_size=12, color=Palette.muted
        )
    if page_no is not None:
        add_textbox(
            slide,
            12.05,
            0.42,
            0.6,
            0.28,
            f"{page_no:02d}",
            font_size=12,
            bold=True,
            color=Palette.accent_dark,
            align=PP_ALIGN.RIGHT,
        )
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.65),
        Inches(1.86),
        Inches(12.0),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = Palette.line
    line.line.color.rgb = Palette.line


def add_footer(slide, text: str = "이천시청·장애인 협회 협력 제안") -> None:
    add_textbox(slide, 0.7, 7.02, 5.8, 0.22, text, font_size=10, color=Palette.muted)


def add_picture_safe(
    slide, name: str, left: float, top: float, width: float, height: float
) -> None:
    path = SCREEN_DIR / name
    if path.exists():
        slide.shapes.add_picture(
            str(path),
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height),
        )


def add_card_title(
    slide,
    left: float,
    top: float,
    width: float,
    title: str,
    body: str,
    *,
    fill: RGBColor,
) -> None:
    add_round_box(slide, left, top, width, 1.55, fill_color=fill)
    add_textbox(
        slide,
        left + 0.18,
        top + 0.14,
        width - 0.36,
        0.3,
        title,
        font_size=15,
        bold=True,
        color=Palette.ink,
    )
    add_textbox(
        slide,
        left + 0.18,
        top + 0.5,
        width - 0.36,
        0.84,
        body,
        font_size=12,
        color=Palette.muted,
    )


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = []
    for _ in range(12):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, Palette.bg)
        add_soft_motif(slide)
        slides.append(slide)

    # 1 Cover
    s = slides[0]
    pill = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(0.62),
        Inches(2.2),
        Inches(0.42),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = Palette.peach_soft
    pill.line.color.rgb = Palette.peach_soft
    add_textbox(
        s,
        0.82,
        0.72,
        2.0,
        0.2,
        "2026 시범사업 제안서",
        font_size=12,
        bold=True,
        color=Palette.accent_dark,
    )
    add_textbox(
        s,
        0.72,
        1.22,
        7.3,
        1.2,
        "꿈꾸는 카메라 기반\nAI 사진·글·음악 창작지원 사업",
        font_size=28,
        bold=True,
        color=Palette.ink,
        font_name="Malgun Gothic",
    )
    add_textbox(
        s,
        0.75,
        2.55,
        5.8,
        0.7,
        "이천시청 · 장애인 협회 제안용\nStory Lens 기반 디지털 문화복지 시범 프로젝트",
        font_size=16,
        color=Palette.muted,
    )
    add_round_box(s, 8.1, 0.9, 4.4, 5.8, fill_color=Palette.surface)
    add_picture_safe(s, "photobook.png", 8.3, 1.15, 4.0, 5.1)
    add_textbox(
        s,
        0.78,
        4.2,
        5.8,
        0.35,
        "핵심 제안",
        font_size=14,
        bold=True,
        color=Palette.accent_dark,
    )
    add_bullets(
        s,
        0.82,
        4.55,
        5.8,
        1.9,
        [
            "사진 편집, AI 글쓰기, AI 음악, 사진집 PDF를 하나의 흐름으로 연결",
            "꿈꾸는 카메라 8명 대상 · 10개월 운영형 시범사업(안)",
            "결과물이 남는 문화예술·디지털 포용 프로그램",
        ],
        font_size=15,
    )
    add_footer(s, "Story Lens · Public Proposal Deck")

    # 2 why now
    s = slides[1]
    add_header(
        s,
        "왜 지금 필요한가",
        "꿈꾸는 카메라에서 먼저 검증하고 다른 프로그램으로 확장할 수 있는 디지털 창작 모델",
    )
    add_round_box(s, 0.72, 2.2, 5.3, 3.95, fill_color=Palette.surface)
    add_textbox(
        s,
        0.95,
        2.45,
        4.8,
        0.3,
        "현장의 간극",
        font_size=16,
        bold=True,
        color=Palette.ink,
    )
    add_bullets(
        s,
        0.98,
        2.88,
        4.75,
        2.9,
        [
            "사진 촬영 활동은 많지만 편집 이후의 기록과 결과물화가 약합니다.",
            "장애 당사자가 스스로 완성 경험을 느끼기보다는 지도자가 마무리를 대신하는 경우가 많습니다.",
            "디지털 도구는 복잡하고, 사진·글·음악·출력물이 분절되어 있습니다.",
        ],
    )
    add_card_title(
        s,
        6.35,
        2.2,
        2.05,
        "자기표현",
        "사진을 통해 감정과 기억을 말하고 싶은 수요",
        fill=Palette.peach_soft,
    )
    add_card_title(
        s,
        8.58,
        2.2,
        2.05,
        "디지털 포용",
        "AI를 보조도구로 활용한 접근성 있는 창작 경험",
        fill=Palette.sage_soft,
    )
    add_card_title(
        s,
        10.81,
        2.2,
        1.8,
        "성과 축적",
        "결과집과 보고자료로 남는 프로그램 구조",
        fill=Palette.blue_soft,
    )
    add_card_title(
        s,
        6.35,
        4.2,
        3.05,
        "행정 친화성",
        "사진집 PDF, 전시 자료, 보고서로 바로 전환 가능한 산출물",
        fill=Palette.surface,
    )
    add_card_title(
        s,
        9.6,
        4.2,
        3.01,
        "지역 확장성",
        "이천시 복지관·학교·협회 협력으로 확산 가능한 구조",
        fill=Palette.surface,
    )
    add_footer(s)

    # 3 solution flow
    s = slides[2]
    add_header(
        s,
        "Story Lens가 제공하는 해결 구조",
        "한 장의 사진이 기록과 창작 결과로 이어지는 6단계 흐름",
    )
    flow = [
        ("① 사진 준비", "촬영 또는 앨범 업로드"),
        ("② 주제 선택", "추천 키워드 또는 직접 입력"),
        ("③ 간편 편집", "필터·슬라이더·회전 조정"),
        ("④ AI 글쓰기", "한 줄 생각 → 최대 5줄 초안"),
        ("⑤ AI 음악", "스타일 선택 후 음악 생성"),
        ("⑥ 사진집 제작", "보관함 축적 후 PDF 결과집"),
    ]
    for idx, (title, desc) in enumerate(flow):
        left = 0.8 + (idx % 3) * 4.1
        top = 2.25 + (idx // 3) * 1.95
        fill = [Palette.peach_soft, Palette.sage_soft, Palette.blue_soft][idx % 3]
        add_round_box(s, left, top, 3.55, 1.45, fill_color=fill)
        add_textbox(
            s,
            left + 0.18,
            top + 0.18,
            3.1,
            0.3,
            title,
            font_size=15,
            bold=True,
            color=Palette.ink,
        )
        add_textbox(
            s,
            left + 0.18,
            top + 0.6,
            3.05,
            0.55,
            desc,
            font_size=12,
            color=Palette.muted,
        )
    add_footer(s)

    # 4 screenshots
    s = slides[3]
    add_header(
        s,
        "시연용 앱 화면",
        "앱의 주요 흐름을 설명하기 위한 시연 화면 기준으로 편집·글쓰기·음악·사진집 구조를 보여줍니다.",
    )
    add_picture_safe(s, "home.png", 0.78, 2.15, 3.0, 2.15)
    add_picture_safe(s, "gallery-detail.png", 3.98, 2.15, 3.0, 2.15)
    add_picture_safe(s, "write.png", 7.18, 2.15, 2.95, 2.15)
    add_picture_safe(s, "music.png", 10.3, 2.15, 2.28, 2.15)
    labels = [
        (0.8, "홈 화면"),
        (4.0, "보관함 상세"),
        (7.2, "AI 글쓰기"),
        (10.3, "AI 음악 생성"),
    ]
    for left, label in labels:
        add_textbox(
            s,
            left,
            4.42,
            2.9,
            0.25,
            label,
            font_size=12,
            bold=True,
            color=Palette.accent_dark,
            align=PP_ALIGN.CENTER,
        )
    add_round_box(s, 0.8, 5.0, 11.8, 1.25, fill_color=Palette.surface)
    add_textbox(
        s,
        1.05,
        5.25,
        11.2,
        0.7,
        "※ 실제 앱은 로그인, 촬영/업로드, 편집, AI 글쓰기, AI 음악, 보관함, 사진집 PDF 흐름이 구현되어 있으며, AI 생성 시간은 외부 서비스 상태에 따라 달라질 수 있습니다.",
        font_size=12,
        color=Palette.muted,
    )
    add_footer(s)

    # 5 beneficiaries
    s = slides[4]
    add_header(
        s,
        "대상자 중심 설계",
        "아동과 성인 모두가 이해하기 쉬운 단계형 UI를 기본 전제로 합니다.",
    )
    add_card_title(
        s,
        0.82,
        2.2,
        3.85,
        "장애아동 참여자",
        "큰 버튼, 직관적인 화면 이동, 사진 한 장에서 출발하는 창작 경험을 통해 자신감을 높입니다.",
        fill=Palette.peach_soft,
    )
    add_card_title(
        s,
        4.78,
        2.2,
        3.85,
        "장애성인 참여자",
        "자신의 기억과 감정을 스스로 표현하고, 사진·글·음악을 연결하는 성취 경험을 제공합니다.",
        fill=Palette.sage_soft,
    )
    add_card_title(
        s,
        8.74,
        2.2,
        3.75,
        "보호자·운영자",
        "결과물을 함께 확인하고 다음 활동으로 이어가기 쉬운 구조를 제공합니다.",
        fill=Palette.blue_soft,
    )
    add_round_box(s, 0.82, 4.25, 11.67, 1.95, fill_color=Palette.surface)
    add_textbox(
        s,
        1.05,
        4.52,
        2.2,
        0.25,
        "운영 시 고려 포인트",
        font_size=14,
        bold=True,
        color=Palette.accent_dark,
    )
    add_bullets(
        s,
        1.05,
        4.9,
        11.0,
        1.05,
        [
            "접근성을 고려한 단계형 UI를 활용하되, 실제 현장 적용 시 참여자별 지원 강도는 다르게 설계합니다.",
            "AI 글쓰기와 음악 생성은 외부 서비스 연동 기능이므로, 수업 중에는 대기 시간을 분리해 운영하는 것이 안정적입니다.",
            "보조인력, 보호자, 활동지원사와의 협력 구조를 함께 설계하면 몰입도와 완성도가 높아집니다.",
        ],
        font_size=13,
    )
    add_footer(s)

    # 6 roadmap
    s = slides[5]
    add_header(
        s,
        "10개월 운영 로드맵(안)",
        "이천시 지역 기반 문화복지 프로그램으로 운영 가능한 월별 구조",
    )
    months = [
        ("3월", "오리엔테이션", Palette.peach_soft),
        ("4월", "봄꽃·감정 기록", Palette.sage_soft),
        ("5월", "풍경과 장소", Palette.blue_soft),
        ("6월", "자연 표현", Palette.peach_soft),
        ("7월", "여름과 움직임", Palette.sage_soft),
        ("8월", "일상 회고", Palette.blue_soft),
        ("9월", "캠핑·야외활동", Palette.peach_soft),
        ("10월", "전통과 문화", Palette.sage_soft),
        ("11월", "AI 창작 집중", Palette.blue_soft),
        ("12월", "회고·사진집", Palette.peach_soft),
    ]
    for idx, (month, title, fill) in enumerate(months):
        left = 0.82 + (idx % 5) * 2.42
        top = 2.2 + (idx // 5) * 1.9
        add_round_box(s, left, top, 2.15, 1.45, fill_color=fill)
        add_textbox(
            s,
            left + 0.16,
            top + 0.16,
            0.8,
            0.25,
            month,
            font_size=14,
            bold=True,
            color=Palette.accent_dark,
        )
        add_textbox(
            s,
            left + 0.16,
            top + 0.53,
            1.75,
            0.55,
            title,
            font_size=12,
            color=Palette.ink,
            bold=True,
        )
    add_round_box(s, 0.82, 5.98, 11.67, 0.42, fill_color=Palette.surface)
    add_textbox(
        s,
        1.0,
        6.08,
        11.2,
        0.2,
        "월별 일정 화면은 운영 관리 기능이 아니라 예시형 활동 콘텐츠로 활용하는 것을 권장합니다.",
        font_size=11,
        color=Palette.muted,
    )
    add_footer(s)

    # 7 session design
    s = slides[6]
    add_header(
        s,
        "회차별 수업 구성 예시",
        "사진 한 장에서 시작해 결과물이 남도록 구성한 100분 프로그램 예시",
    )
    stages = [
        ("10분", "도입", "오늘의 주제 소개 · 감정 단어 열기"),
        ("25분", "촬영", "현장 또는 앨범 이미지 탐색"),
        ("25분", "편집", "필터·슬라이더·회전 조정"),
        ("20분", "글쓰기", "한 줄 생각 입력 · AI 초안 보조"),
        ("20분", "확장", "AI 음악 또는 결과 공유"),
    ]
    for idx, (time, title, body) in enumerate(stages):
        left = 0.82 + idx * 2.48
        add_round_box(
            s,
            left,
            2.45,
            2.1,
            2.3,
            fill_color=[
                Palette.peach_soft,
                Palette.sage_soft,
                Palette.blue_soft,
                Palette.peach_soft,
                Palette.sage_soft,
            ][idx],
        )
        add_textbox(
            s,
            left + 0.15,
            2.65,
            0.9,
            0.22,
            time,
            font_size=12,
            bold=True,
            color=Palette.accent_dark,
        )
        add_textbox(
            s,
            left + 0.15,
            3.05,
            1.6,
            0.24,
            title,
            font_size=14,
            bold=True,
            color=Palette.ink,
        )
        add_textbox(
            s, left + 0.15, 3.45, 1.7, 0.8, body, font_size=11, color=Palette.muted
        )
    add_picture_safe(s, "gallery-detail.png", 8.65, 5.0, 3.45, 1.95)
    add_textbox(
        s,
        0.86,
        5.15,
        7.2,
        0.8,
        "운영자는 보관함 상세 화면을 통해 사진·글·음악을 함께 확인하며 참여자의 변화와 성취를 점검할 수 있습니다.",
        font_size=14,
        color=Palette.muted,
    )
    add_footer(s)

    # 8 impact
    s = slides[7]
    add_header(
        s, "기대 효과와 산출물", "참여자, 가족, 기관 모두에게 남는 결과를 만듭니다."
    )
    add_card_title(
        s,
        0.82,
        2.2,
        3.65,
        "참여자 변화",
        "사진을 통해 감정을 표현하고, 글과 음악으로 확장하는 자기표현 경험",
        fill=Palette.peach_soft,
    )
    add_card_title(
        s,
        4.86,
        2.2,
        3.65,
        "가족과 보호자",
        "눈으로 확인 가능한 결과물과 사진집이 가족 공유 자료가 됨",
        fill=Palette.sage_soft,
    )
    add_card_title(
        s,
        8.9,
        2.2,
        3.65,
        "기관 성과",
        "전시·보고서·발표 자료로 활용할 수 있는 시각적 성과 확보",
        fill=Palette.blue_soft,
    )
    add_round_box(s, 0.82, 4.35, 6.2, 1.85, fill_color=Palette.surface)
    add_textbox(
        s,
        1.05,
        4.58,
        2.0,
        0.25,
        "대표 산출물(안)",
        font_size=14,
        bold=True,
        color=Palette.accent_dark,
    )
    add_bullets(
        s,
        1.05,
        4.95,
        5.6,
        0.95,
        [
            "사진·글·음악 연계 결과물 48건 목표",
            "개인 또는 공동 사진집 PDF 1종 이상",
            "성과공유회 또는 전시용 인쇄물 제작",
        ],
        font_size=13,
    )
    add_picture_safe(s, "home.png", 7.35, 4.33, 5.18, 1.95)
    add_footer(s)

    # 9 KPI
    s = slides[8]
    add_header(s, "성과지표(KPI) 제안", "시범사업 평가에 바로 연결할 수 있는 기본 지표")
    metrics = [
        ("직접 참여", "8명", "장애아동 및 성인"),
        ("운영 협력", "7명", "선생님 6 + 복지사 1"),
        ("1인당 결과물", "6건+", "사진 또는 연계 창작물"),
        ("긍정 응답", "80%+", "참여자·보호자 의견"),
    ]
    for idx, (title, value, desc) in enumerate(metrics):
        left = 0.85 + idx * 3.05
        add_round_box(
            s,
            left,
            2.35,
            2.65,
            2.1,
            fill_color=[
                Palette.peach_soft,
                Palette.sage_soft,
                Palette.blue_soft,
                Palette.peach_soft,
            ][idx],
        )
        add_textbox(
            s,
            left + 0.16,
            2.6,
            2.2,
            0.22,
            title,
            font_size=13,
            bold=True,
            color=Palette.accent_dark,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            s,
            left + 0.16,
            3.0,
            2.2,
            0.55,
            value,
            font_size=28,
            bold=True,
            color=Palette.ink,
            align=PP_ALIGN.CENTER,
            font_name="Malgun Gothic",
        )
        add_textbox(
            s,
            left + 0.16,
            3.7,
            2.2,
            0.3,
            desc,
            font_size=11,
            color=Palette.muted,
            align=PP_ALIGN.CENTER,
        )
    add_round_box(s, 0.85, 4.9, 11.6, 1.35, fill_color=Palette.surface)
    add_textbox(
        s,
        1.05,
        5.15,
        11.1,
        0.72,
        "※ 본 지표는 시범사업 제안 기준이며, 실제 운영 시 기관 목표와 예산 조건에 따라 조정 가능합니다. AI 생성 기능은 외부 서비스 상태에 따라 시간이 달라질 수 있으므로, 성공률과 소요 시간도 함께 기록하는 것을 권장합니다.",
        font_size=12,
        color=Palette.muted,
    )
    add_footer(s)

    # 10 budget summary
    s = slides[9]
    add_header(
        s,
        "예산 총괄(추정안)",
        "총 사업비 14,000,000원 기준 · 꿈꾸는 카메라 1차 적용 가정",
    )
    add_round_box(s, 0.82, 2.2, 3.2, 1.9, fill_color=Palette.peach_soft)
    add_textbox(
        s,
        1.05,
        2.52,
        2.4,
        0.25,
        "총 사업비(추정)",
        font_size=14,
        bold=True,
        color=Palette.accent_dark,
    )
    add_textbox(
        s,
        1.0,
        3.0,
        2.7,
        0.45,
        "14,000,000원",
        font_size=28,
        bold=True,
        color=Palette.ink,
        font_name="Malgun Gothic",
        align=PP_ALIGN.CENTER,
    )
    categories = [
        ("운영기획", 3.0, Palette.peach_soft),
        ("프로그램 운영", 2.4, Palette.sage_soft),
        ("결과물 제작", 2.8, Palette.blue_soft),
        ("플랫폼·AI", 2.08, Palette.surface),
        ("확장·예비비", 3.72, Palette.peach_soft),
    ]
    x = 4.45
    for title, value, fill in categories:
        width = (value / 3.72) * 2.1
        bar = s.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(2.45),
            Inches(width),
            Inches(0.45),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = fill
        bar.line.color.rgb = fill
        add_textbox(
            s,
            x,
            3.02,
            max(width, 1.2),
            0.2,
            f"{title} {value:.1f}백만원",
            font_size=11,
            color=Palette.muted,
        )
        x += width + 0.12
    add_round_box(s, 4.45, 3.7, 8.0, 2.1, fill_color=Palette.surface)
    add_bullets(
        s,
        4.72,
        4.02,
        7.4,
        1.45,
        [
            "이천시청 지원금 10,000,000원(안)",
            "협회/운영기관 분담 2,500,000원(안)",
            "후원·협찬 연계 1,500,000원(안)",
            "서버 100만원/년, AI 약 9만원/월 기준이며 확장 시 비용 증가 가능",
        ],
        font_size=14,
    )
    add_footer(s)

    # 11 roles
    s = slides[10]
    add_header(
        s, "협력 구조와 역할 분담", "공공재원, 협회 운영, 실무팀 실행이 나뉘는 구조"
    )
    add_card_title(
        s,
        0.82,
        2.15,
        3.8,
        "이천시청",
        "사업비 지원, 공공 협력, 지역 확산, 성과 보고 연계",
        fill=Palette.peach_soft,
    )
    add_card_title(
        s,
        4.77,
        2.15,
        3.8,
        "장애인 협회 및 복지기관",
        "참여자 모집, 현장 운영, 보호자 소통, 일정 협력",
        fill=Palette.sage_soft,
    )
    add_card_title(
        s,
        8.72,
        2.15,
        3.8,
        "운영팀",
        "교육 설계, 앱 운영, 강사 배치, 결과집 제작",
        fill=Palette.blue_soft,
    )
    add_round_box(s, 1.1, 4.55, 11.0, 1.45, fill_color=Palette.surface)
    add_textbox(
        s,
        1.35,
        4.82,
        10.4,
        0.52,
        "권장 운영 형태: 꿈꾸는 카메라 현장 운영 + 협회 협력 + 시청 지원 + 향후 프로그램 확장",
        font_size=17,
        bold=True,
        color=Palette.ink,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        1.35,
        5.35,
        10.4,
        0.3,
        "→ 공공성, 지속성, 참여자 지원 강도를 함께 확보할 수 있는 구조",
        font_size=12,
        color=Palette.muted,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s)

    # 12 closing
    s = slides[11]
    add_header(
        s, "제안 결론 및 요청 사항", "사진 한 장에서 시작하는 디지털 문화복지 시범사업"
    )
    add_round_box(s, 0.82, 2.15, 7.25, 3.55, fill_color=Palette.surface)
    add_textbox(
        s,
        1.05,
        2.45,
        6.7,
        0.4,
        "왜 Story Lens인가",
        font_size=17,
        bold=True,
        color=Palette.accent_dark,
    )
    add_bullets(
        s,
        1.08,
        2.95,
        6.55,
        2.2,
        [
            "사진 편집, AI 글쓰기, AI 음악, 사진집 제작까지 실제 구현된 흐름을 기반으로 함",
            "장애아동과 성인이 자기표현 결과를 직접 완성하는 경험을 중심에 둠",
            "성과물이 눈에 보이기 때문에 공공기관 보고와 지역 확산에 적합함",
            "시범사업 운영 후 학교·복지관·축제 연계로 확장 가능한 구조",
        ],
        font_size=15,
    )
    add_round_box(s, 8.35, 2.15, 4.15, 3.55, fill_color=Palette.peach_soft)
    add_textbox(
        s,
        8.65,
        2.45,
        3.5,
        0.35,
        "협의 요청",
        font_size=17,
        bold=True,
        color=Palette.accent_dark,
    )
    add_bullets(
        s,
        8.65,
        2.95,
        3.45,
        1.8,
        [
            "꿈꾸는 카메라 1차 적용 예산 구조 협의",
            "운영기관/협회 역할 분담 확정",
            "참여자 8명 기준 운영안 확정",
            "확장 시 웹 운영비·AI 비용 증가 구조 논의",
        ],
        font_size=14,
    )
    add_textbox(
        s,
        0.85,
        6.1,
        11.6,
        0.5,
        "제안서, 예산안, 앱 매뉴얼, 실제 앱 스크린샷을 함께 제시하여 협의용 설명 자료로 바로 활용할 수 있도록 구성했습니다.",
        font_size=14,
        color=Palette.muted,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        s,
        "Story Lens · pilot proposal deck for Icheon City Hall & disability associations",
    )

    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_deck()
    print(path)
